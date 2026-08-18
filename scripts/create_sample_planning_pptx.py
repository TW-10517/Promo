"""
Generate a new sample Planning Document PowerPoint (.pptx)
for Dospara / raytrek AI Creator PC Promotion.
"""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

root_dir = Path(__file__).resolve().parent.parent

# Theme Colors
DARK_BLUE = RGBColor(0x1B, 0x36, 0x5D)
ACCENT_TEAL = RGBColor(0x00, 0x82, 0x8A)
TEXT_GRAY = RGBColor(0x33, 0x33, 0x33)


def add_slide_header(slide, title_text: str, category_text: str = "RAYTREK PROMOTION PLAN 2026"):
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True

    p_cat = tf.paragraphs[0]
    p_cat.text = category_text
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_TEAL

    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = DARK_BLUE


def create_sample_presentation(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 Widescreen

    blank_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # SLIDE 1: Cover Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    box1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(8.0), Inches(3.2))
    tf1 = box1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "【raytrek / ドスパラ】"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL

    p2 = tf1.add_paragraph()
    p2.text = "2026年秋 最新AIクリエイター向けPC\n「raytrek AI Studio」発売記念プロモーション企画書"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = DARK_BLUE

    p3 = tf1.add_paragraph()
    p3.text = "\n起案日: 2026年9月1日（火）\n起案部門: 販売促進部 プロモーション企画課\n担当者: 佐藤 健一 (企画担当: 高橋 翼)\n対象ブランド: raytrek / THIRDWAVE"
    p3.font.size = Pt(12)
    p3.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 2: Background & Purpose
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide2, "1. 企画の背景と目的")

    box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf2 = box2.text_frame
    tf2.word_wrap = True

    tf2.paragraphs[0].text = "■ 開発・市場背景"
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = DARK_BLUE

    p_bg = tf2.add_paragraph()
    p_bg.text = (
        "・生成AI（画像・動画・3D・音楽）技術の急速な進化に伴い、イラストレーター、3Dモデラー、動画クリエイターにおける「ローカルAI環境」の需要が急拡大。\n"
        "・VRAM容量（16GB〜32GB以上）と大容量高速メモリ（64GB〜128GB）を搭載したハイエンド制作環境への乗り換え需要が顕在化している。"
    )
    p_bg.font.size = Pt(12)
    p_bg.font.color.rgb = TEXT_GRAY

    p_pur_head = tf2.add_paragraph()
    p_pur_head.text = "\n■ 本プロモーションの目的"
    p_pur_head.font.size = Pt(14)
    p_pur_head.font.bold = True
    p_pur_head.font.color.rgb = DARK_BLUE

    p_pur = tf2.add_paragraph()
    p_pur.text = (
        "・新世代AIクリエイター向けPC「raytrek AI Studio」の発売に合わせて認知拡大と初期販売数を最大化する。\n"
        "・購入サポートクーポンやメモリ増量キャンペーンを通じて、制作現場やフリーランスの買い替えハードルを低減し、クリエイターPC市場におけるraytrekのプレゼンスを強化する。"
    )
    p_pur.font.size = Pt(12)
    p_pur.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 3: Target Audience & Positioning
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide3, "2. ターゲット層 & ポジショニング")

    box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf3 = box3.text_frame
    tf3.word_wrap = True

    tf3.paragraphs[0].text = "■ メインターゲット層"
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = DARK_BLUE

    p_target = tf3.add_paragraph()
    p_target.text = (
        "1. プロCGクリエイター / 3Dデザイナー (20代〜50代)\n"
        "   - Maya, Blender, Unreal Engine 5および生成AIを活用した超高負荷レンダリングを日常的に行う制作会社・フリーランス\n"
        "2. イラストレーター / アニメーション制作者\n"
        "   - CLIP STUDIO PAINT、Photoshopでの超高解像度イラスト制作やローカルAI画像生成の支援機能を活用する層\n"
        "3. 動画編集者 / VFXアーティスト\n"
        "   - Premiere Pro, DaVinci Resolveによる8K動画編集やエフェクト処理を高速に行いたい層"
    )
    p_target.font.size = Pt(12)
    p_target.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 4: Promotion Overview
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide4, "3. 施策概要 & キャンペーン骨子")

    box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf4 = box4.text_frame
    tf4.word_wrap = True

    tf4.paragraphs[0].text = "■ キャンペーン概要（raytrek AIクリエイター応援フェア）"
    tf4.paragraphs[0].font.size = Pt(14)
    tf4.paragraphs[0].font.bold = True
    tf4.paragraphs[0].font.color.rgb = DARK_BLUE

    p_promo = tf4.add_paragraph()
    p_promo.text = (
        "1. 最大80,000円分 raytrek購入サポートクーポン配布\n"
        "   - 対象のraytrekデスクトップ・ノートPC購入時に利用可能な即時値引きクーポンを提供\n"
        "2. メモリ無償アップグレード（64GB → 128GB）キャンペーン\n"
        "   - フラッグシップモデル購入者を対象に、メモリ倍増カスタマイズを無料提供（先着500名）\n"
        "3. Web特設ランディングページの開設 & 秋葉原店舗クリエイター体験コーナー設置\n"
        "   - 特設ページ: https://www.dospara.co.jp/event/raytrek-ai-studio-2026.html\n"
        "   - 全国ドスパラ主要店舗にてAdobe Creative Cloud & Stable Diffusion実機動作デモを実施"
    )
    p_promo.font.size = Pt(12)
    p_promo.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 5: Target Products
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide5, "4. 主要対象製品ラインナップ")

    box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf5 = box5.text_frame
    tf5.word_wrap = True

    tf5.paragraphs[0].text = "■ 主要モデルスペック・価格一例"
    tf5.paragraphs[0].font.size = Pt(14)
    tf5.paragraphs[0].font.bold = True
    tf5.paragraphs[0].font.color.rgb = DARK_BLUE

    p_prods = tf5.add_paragraph()
    p_prods.text = (
        "● raytrek 4C-X9 (フラッグシップAIワークステーション)\n"
        "   - 構成: Intel Core Ultra 9 285K / GeForce RTX 5090 32GB / メモリ128GB DDR5 / SSD 4TB NVMe Gen5\n"
        "   - 特徴: ローカルLLM検証から8K動画編集、リアルタイム3D生成まで極限の処理性能を提供。クーポン利用で80,000円購入サポート\n"
        "   - 販売URL: https://www.dospara.co.jp/TC30/MC31090.html\n\n"
        "● raytrek 4C-T7 (スタンダードクリエイターデスクトップ)\n"
        "   - 構成: AMD Ryzen 9 9900X / GeForce RTX 5070Ti 16GB / メモリ64GB / SSD 2TB NVMe Gen4\n"
        "   - 特徴: イラスト・3DCG・配信に最適な高コストパフォーマンス機。クーポン利用で40,000円購入サポート\n"
        "   - 販売URL: https://www.dospara.co.jp/TC30/MC31070.html\n\n"
        "● raytrek R6-AA (16インチ ハイエンドクリエイターノート)\n"
        "   - 構成: Core Ultra 7 265H / RTX 5070 Laptop 8GB / 16.0型 4K有機EL (DCI-P3 100%) / メモリ32GB。20,000円購入サポート"
    )
    p_prods.font.size = Pt(11)
    p_prods.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 6: Channels, Schedule & Budget
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide6, "5. 展開チャネル・スケジュール・予算")

    box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf6 = box6.text_frame
    tf6.word_wrap = True

    tf6.paragraphs[0].text = "■ 展開チャネル・期間・費用"
    tf6.paragraphs[0].font.size = Pt(14)
    tf6.paragraphs[0].font.bold = True
    tf6.paragraphs[0].font.color.rgb = DARK_BLUE

    p_plan = tf6.add_paragraph()
    p_plan.text = (
        "・販売チャネル: 全国のドスパラ各店舗、およびドスパラ公式通販サイト (Web & B2B法人窓口)\n"
        "・実施期間: 2026年9月15日(火) 11:00 〜 2026年10月15日(木) 10:59 (店舗は14日閉店まで)\n"
        "・スケジュール進行:\n"
        "   - 9月15日(火) 11:00 公式プレスリリース配信・新製品発売・特設ページオープン\n"
        "   - 9月19日(土) ドスパラ秋葉原本店にて著名イラストレーター登壇AI活用セミナー開催\n"
        "   - 10月15日(木) 10:59 Webプロモーション終了\n"
        "・概算予算: プロモーション全体予算 2,800万円 (Web広告、特設ページ制作、店舗デモ機設置、クーポン原資)\n"
        "・期待効果: 期間中raytrekシリーズ販売台数 前年同期比 140%達成、クリエイター新規会員登録数 5,000名獲得"
    )
    p_plan.font.size = Pt(11)
    p_plan.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 7: Risks & Considerations
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide7, "6. 想定リスク・留意事項 & 問い合わせ窓口")

    box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf7 = box7.text_frame
    tf7.word_wrap = True

    tf7.paragraphs[0].text = "■ 想定リスク及び対策"
    tf7.paragraphs[0].font.size = Pt(14)
    tf7.paragraphs[0].font.bold = True
    tf7.paragraphs[0].font.color.rgb = DARK_BLUE

    p_risk = tf7.add_paragraph()
    p_risk.text = (
        "1. DDR5 128GB大容量メモリの調達リードタイム長期化\n"
        "   → 発売開始前に先行して1,000セット分のメモリ在庫を国内工場に確保済み\n"
        "2. クリエイターソフトウェア互換性に関する問い合わせ増加\n"
        "   → 主要クリエイティブソフト（Adobe, Blender, DaVinci等）動作検証済みリストを特設サイトに事前掲載\n\n"
        "■ 備考・連絡先\n"
        "・広報問い合わせ先: 株式会社サードウェーブ 広報室 (TEL: 03-5294-2043 / dospara-koho@twave.co.jp)\n"
        "・サポートセンター: 03-4332-9193 (24時間365日対応)"
    )
    p_risk.font.size = Pt(11)
    p_risk.font.color.rgb = TEXT_GRAY

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"[OK] Created PowerPoint at: {output_path}")


if __name__ == "__main__":
    out = root_dir / "data" / "raytrek_ai_pc_planning.pptx"
    create_sample_presentation(out)
