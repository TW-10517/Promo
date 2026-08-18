"""
Creates a professional demo Planning Document PowerPoint (.pptx)
for Dospara / GALLERIA promotions.
"""

from __future__ import annotations

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

root_dir = Path(__file__).resolve().parent.parent

# Colors
DARK_BLUE = RGBColor(0x1B, 0x36, 0x5D)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
TEXT_GRAY = RGBColor(0x33, 0x33, 0x33)
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)


def add_slide_header(slide, title_text: str, category_text: str = "GALLERIA PROMOTION PLAN 2026"):
    # Header title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_cat = tf.paragraphs[0]
    p_cat.text = category_text
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_BLUE
    
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = DARK_BLUE


def create_demo_presentation(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625) # 16:9 Widescreen

    blank_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # SLIDE 1: Cover Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    
    box1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(8.0), Inches(3.2))
    tf1 = box1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "【GALLERIA / ドスパラ】"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    p2 = tf1.add_paragraph()
    p2.text = "2026年秋 最新AIゲーミングPC\n発売記念プロモーション企画書"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = DARK_BLUE
    
    p3 = tf1.add_paragraph()
    p3.text = "\n起案日: 2026年8月14日（金）\n起案部門: 販売促進部 プロモーション企画課\n担当者: 永井 正樹 (企画担当: 山田 太郎)\n対象ブランド: GALLERIA / THIRDWAVE"
    p3.font.size = Pt(12)
    p3.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 2: Background & Purpose (背景と目的)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide2, "1. 企画の背景と目的")
    
    box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    
    tf2.paragraphs[0].text = "■ 開発・市場背景"
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = DARK_BLUE
    
    p_bg = tf2.add_paragraph()
    p_bg.text = (
        "・最新世代CPU（Ryzen 7 9800X3D / Core Ultra）および次世代GPU（GeForce RTX 50シリーズ）の登場により、"
        "PCゲーマーや動画クリエイターの買い替え需要が急速に高まっている。\n"
        "・特に4K/高フレームレートゲームやローカルAI画像生成・配信の負荷増大に対応するハイエンドPCへの注目度が顕著に上昇。"
    )
    p_bg.font.size = Pt(12)
    p_bg.font.color.rgb = TEXT_GRAY
    
    p_pur_head = tf2.add_paragraph()
    p_pur_head.text = "\n■ 本プロモーションの目的"
    p_pur_head.font.size = Pt(14)
    p_pur_head.font.bold = True
    p_pur_head.font.color.rgb = DARK_BLUE
    
    p_pur = tf2.add_paragraph()
    p_pur.text = (
        "・秋の需要期に合わせた強力な購入サポートとアップグレード施策により、GALLERIAフラッグシップモデルの認知拡大と初期販売数を最大化する。\n"
        "・「圧倒的パフォーマンスと安心の国内生産・サポート」を訴求し、ゲーミングPC市場におけるドスパラのブランドリーダーシップを確立する。"
    )
    p_pur.font.size = Pt(12)
    p_pur.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 3: Target Audience & Positioning (対象顧客)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide3, "2. ターゲット層 & ポジショニング")
    
    box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    
    tf3.paragraphs[0].text = "■ メインターゲット層"
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = DARK_BLUE
    
    p_target = tf3.add_paragraph()
    p_target.text = (
        "1. コアゲーマー (20代〜40代 男性・女性)\n"
        "   - 最新3Aタイトルや競技性の高いFPS/格闘ゲームで最高画質・高リフレッシュレート環境を求める層\n"
        "2. ゲーム配信者 / VTuber / 動画クリエイター\n"
        "   - 3Dアバター制御、OBS高ビットレート配信、4K動画レンダリングを1台で完結させたい層\n"
        "3. ビジネス / AI開発従事者・学生\n"
        "   - ローカルLLM検証や画像生成AIを日常業務・研究で活用するパワーユーザー"
    )
    p_target.font.size = Pt(12)
    p_target.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 4: Promotion Overview & Campaign (施策概要)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide4, "3. 施策概要 & キャンペーン骨子")
    
    box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf4 = box4.text_frame
    tf4.word_wrap = True
    
    tf4.paragraphs[0].text = "■ キャンペーン概要（夏・秋の大感謝SALE連動）"
    tf4.paragraphs[0].font.size = Pt(14)
    tf4.paragraphs[0].font.bold = True
    tf4.paragraphs[0].font.color.rgb = DARK_BLUE
    
    p_promo = tf4.add_paragraph()
    p_promo.text = (
        "1. 最大145,000円分 新品PC購入サポートクーポン配布\n"
        "   - フラッグシップデスクトップからエントリーノートまで、製品グレードに応じた即時値引きクーポンを提供\n"
        "2. グラフィックボードお得アップグレードキャンペーン\n"
        "   - RTX 5060搭載指定PC購入時、わずか3,000円でRTX 5060Ti（8GB）へアップグレード可能\n"
        "3. Web特設ランディングページの開設 & 店頭連動展開\n"
        "   - 特設ページ: https://www.dospara.co.jp/event/autumn-thanksgiving2026.html\n"
        "   - 秋葉原本店および全国ドスパラ店舗での実機体験デモ展示"
    )
    p_promo.font.size = Pt(12)
    p_promo.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 5: Target Products (対象製品ラインナップ)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide5, "4. 主要対象製品ラインナップ")
    
    box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf5 = box5.text_frame
    tf5.word_wrap = True
    
    tf5.paragraphs[0].text = "■ 主要モデルスペック・価格一例"
    tf5.paragraphs[0].font.size = Pt(14)
    tf5.paragraphs[0].font.bold = True
    tf5.paragraphs[0].font.color.rgb = DARK_BLUE
    
    p_prods = tf5.add_paragraph()
    p_prods.text = (
        "● GALLERIA XDR7A-R58-WL (ハイパフォーマンスデスクトップ)\n"
        "   - 構成: AMD Ryzen 7 9800X3D / GeForce RTX 5080 16GB / メモリ64GB / SSD 2TB Gen4\n"
        "   - 特徴: 4Kウルトラ画質・最高峰ゲームプレイとAI学習に最適。クーポン利用で145,000円購入サポート\n"
        "   - 販売URL: https://www.dospara.co.jp/TC30/MC25181.html\n\n"
        "● GALLERIA NMC9L-R58-H6 (ハイエンドゲーミングノート)\n"
        "   - 構成: Core Ultra 9 290HX Plus / RTX 5080 16GB / 16.0型 300Hz WQXGA / メモリ32GB\n"
        "   - 特徴: 持ち運び可能な本体に最高峰パーツ凝縮。クーポン利用で50,000円購入サポート\n\n"
        "● THIRDWAVE F-14BR5A (スタンダードノート)\n"
        "   - 構成: Ryzen 5 7430U / 16GBメモリ / 500GB SSD / 14.0型 フルHD。1,000円購入サポート"
    )
    p_prods.font.size = Pt(11)
    p_prods.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 6: Channels, Schedule & Budget (チャネル・日程・予算)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide6, "5. 展開チャネル・スケジュール・予算")
    
    box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf6 = box6.text_frame
    tf6.word_wrap = True
    
    tf6.paragraphs[0].text = "■ 展開チャネル・期間・費用"
    tf6.paragraphs[0].font.size = Pt(14)
    tf6.paragraphs[0].font.bold = True
    tf6.paragraphs[0].font.color.rgb = DARK_BLUE
    
    p_plan = tf6.add_paragraph()
    p_plan.text = (
        "・販売チャネル: 全国のドスパラ各店舗、およびドスパラ公式通販サイト\n"
        "・実施期間: 2026年8月14日(金) 11:00 〜 2026年8月28日(金) 10:59 (店舗は27日閉店まで)\n"
        "・スケジュール進行:\n"
        "   - 8月14日(金) 11:00 公式プレスリリース配信・セール開始・特設ページオープン\n"
        "   - 8月15日(土) 店頭体験イベント・SNS告知キャンペーン開始\n"
        "   - 8月28日(金) 10:59 Webセール終了\n"
        "・概算予算: プロモーション全体予算 3,500万円 (Web広告、店頭販促物、クーポン原資含む)\n"
        "・期待効果: 期間中GALLERIA販売台数 前年同期比 135%達成、新規会員登録数 8,000名獲得"
    )
    p_plan.font.size = Pt(11)
    p_plan.font.color.rgb = TEXT_GRAY

    # -------------------------------------------------------------
    # SLIDE 7: Risks & Considerations (想定リスク・留意点)
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide7, "6. 想定リスク・留意事項 & 問い合わせ窓口")
    
    box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
    tf7 = box7.text_frame
    tf7.word_wrap = True
    
    tf7.paragraphs[0].text = "■ 想定リスク及び対策"
    tf7.paragraphs[0].font.size = Pt(14)
    tf7.paragraphs[0].font.bold = True
    tf7.paragraphs[0].font.color.rgb = DARK_BLUE
    
    p_risk = tf7.add_paragraph()
    p_risk.text = (
        "1. RTX 5080等 ハイエンド部材の初期ロット供給不足\n"
        "   → 事前サプライチェーン調整により初期2,000台分の部材を先行確保済み\n"
        "2. セール初日のアクセス集中および店舗混雑\n"
        "   → サーバーオートスケール対応および店舗整理券システムの配備\n\n"
        "■ 備考・連絡先\n"
        "・広報問い合わせ先: 株式会社サードウェーブ 広報室 (TEL: 03-5294-2043 / dospara-koho@twave.co.jp)\n"
        "・サポートセンター: 03-4332-9193 (24時間365日対応)"
    )
    p_risk.font.size = Pt(11)
    p_risk.font.color.rgb = TEXT_GRAY

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"[OK] Created PowerPoint at: {output_path}")


if __name__ == "__main__":
    out = root_dir / "data" / "demo_planning_document.pptx"
    create_demo_presentation(out)
