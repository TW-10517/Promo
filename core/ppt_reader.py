"""
PowerPoint text extractor.

Reads a Planning Document (.pptx) and returns every piece of readable text:
shape text, grouped shapes, table cells, chart titles where available and the
speaker notes. The result is a single plain-text string with slide markers,
which is what gets fed to Claude as the source content.

This module is intentionally dependency-light and can be run standalone:

    python -m core.ppt_reader "path/to/planning.pptx"
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.exc import PackageNotFoundError


class PptReadError(Exception):
    """Raised when the .pptx cannot be opened or contains no readable text."""


def _shape_text(shape) -> list[str]:
    """
    Extract text from a single shape.

    Handles the three cases that actually matter in practice:
    group shapes (recurse), tables (cell by cell) and normal text frames.
    """
    parts: list[str] = []

    # 1. Group shape -> recurse into its children.
    #    shape_type 6 == MSO_SHAPE_TYPE.GROUP
    if getattr(shape, "shape_type", None) == 6:
        for child in shape.shapes:
            parts.extend(_shape_text(child))
        return parts

    # 2. Table -> read every cell, keep the row structure with tab separators.
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
        return parts

    # 3. Chart -> the chart title is often the only readable text.
    if getattr(shape, "has_chart", False):
        try:
            chart = shape.chart
            if chart.has_title and chart.chart_title.has_text_frame:
                title = chart.chart_title.text_frame.text.strip()
                if title:
                    parts.append(f"[グラフ] {title}")
        except Exception:  # noqa: BLE001 - malformed charts must not abort extraction
            pass
        return parts

    # 4. Ordinary text frame.
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            line = "".join(run.text for run in paragraph.runs).strip()
            if not line:
                # Some paragraphs carry text directly rather than in runs.
                line = paragraph.text.strip()
            if line:
                parts.append(line)

    return parts


def extract_text(pptx_path: str | Path) -> str:
    """
    Extract all readable text from a .pptx file.

    Parameters
    ----------
    pptx_path:
        Path to the Planning Document.

    Returns
    -------
    str
        Plain text, one block per slide, prefixed with ``--- スライド N ---``.

    Raises
    ------
    PptReadError
        The file is missing, is not a valid .pptx, or holds no readable text.
    """
    path = Path(pptx_path)
    if not path.is_file():
        raise PptReadError(f"ファイルが見つかりません: {path}")

    try:
        presentation = Presentation(str(path))
    except PackageNotFoundError as exc:
        raise PptReadError(
            f"PowerPointファイルとして開けません: {path.name}\n"
            "拡張子が .pptx のファイルを選択してください（.ppt は非対応です）。"
        ) from exc
    except Exception as exc:  # noqa: BLE001 - surface any parser failure to the UI
        raise PptReadError(f"PowerPointファイルの読み込みに失敗しました: {exc}") from exc

    blocks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []

        for shape in slide.shapes:
            lines.extend(_shape_text(shape))

        # Speaker notes frequently hold the detail the slide body omits.
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"[ノート] {notes}")
        except Exception:  # noqa: BLE001 - a broken notes slide must not abort extraction
            pass

        if lines:
            blocks.append(f"--- スライド {index} ---\n" + "\n".join(lines))

    text = "\n\n".join(blocks).strip()
    if not text:
        raise PptReadError(
            f"{path.name} からテキストを抽出できませんでした。\n"
            "スライドが画像のみで構成されている可能性があります。"
        )
    return text


def summarize(text: str) -> str:
    """Short human-readable stat line shown in the UI after extraction."""
    slides = text.count("--- スライド ")
    return f"{slides} スライド / {len(text):,} 文字を抽出しました"


if __name__ == "__main__":  # Standalone verification of step 1 of the build order.
    import sys

    if len(sys.argv) != 2:
        print("使い方: python -m core.ppt_reader <planning.pptx>")
        raise SystemExit(2)
    try:
        extracted = extract_text(sys.argv[1])
    except PptReadError as err:
        print(f"エラー: {err}")
        raise SystemExit(1)
    print(extracted)
    print("\n" + "=" * 60)
    print(summarize(extracted))
